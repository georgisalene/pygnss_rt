#!/usr/bin/env python3
"""
GNSS Quality Monitoring Report Generator

Generates a comprehensive PowerPoint presentation with all quality metrics,
comparisons with CODE products, and detailed analysis.
"""

import os
import sys
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Alias for convenience
RgbColor = RGBColor

# Add parent to path for sibling module imports
sys.path.insert(0, str(Path(__file__).parent))

from pygnss_rt.frontend.db_models import QualityDatabase
from pygnss_rt.frontend.config import DB_PATH
from code_snx_parser import parse_code_snx_file, download_code_snx, get_code_coords_for_station, compute_enu_diff
from code_tro_parser import parse_code_tro_file, download_code_tro, get_code_ztd_for_station

# Output directory for figures
FIG_DIR = Path("/tmp/gnss_report_figures")
FIG_DIR.mkdir(exist_ok=True)

# Colors
COLORS = {
    'primary': '#1a5276',
    'secondary': '#2874a6',
    'accent': '#3498db',
    'success': '#27ae60',
    'warning': '#f39c12',
    'danger': '#e74c3c',
    'gps': '#2ecc71',
    'galileo': '#3498db',
    'glonass': '#e74c3c'
}


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(26, 82, 118)
    shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(40, 116, 166)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_items=None, image_path=None, two_columns=False):
    """Add a content slide with optional image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RgbColor(26, 82, 118)
    title_shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    if image_path and os.path.exists(image_path):
        if content_items:
            # Image on right, text on left
            slide.shapes.add_picture(str(image_path), Inches(5), Inches(1.2), width=Inches(4.8))
            text_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.5), Inches(5.5))
        else:
            # Full width image
            slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.2), width=Inches(9))
    elif content_items:
        text_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5))

    if content_items:
        tf = text_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            if item.startswith("  "):
                p.level = 1

    return slide


def add_table_slide(prs, title, df, max_rows=12):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RgbColor(26, 82, 118)
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Limit rows
    if len(df) > max_rows:
        df = df.head(max_rows)

    rows, cols = df.shape
    rows += 1  # Header row

    # Calculate table dimensions
    table_width = min(9.4, cols * 1.2)
    col_width = table_width / cols

    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1), Inches(table_width), Inches(5.5)).table

    # Header row
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(40, 116, 166)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)

    # Data rows
    for i, (_, row) in enumerate(df.iterrows()):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val) if not pd.isna(val) else ""
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(235, 245, 251)

    return slide


def generate_coordinate_repeatability_figure(db, year, start_doy, end_doy):
    """Generate coordinate repeatability bar chart"""
    stations = db.get_all_stations()

    repeatability_data = []
    for station in stations:
        rep = db.get_coordinate_repeatability(station, year, start_doy, end_doy)
        if rep and rep.get('std_x'):
            repeatability_data.append({
                'Station': station,
                'N (mm)': rep['std_y'] * 1000 if rep['std_y'] else 0,
                'E (mm)': rep['std_x'] * 1000 if rep['std_x'] else 0,
                'U (mm)': rep['std_z'] * 1000 if rep['std_z'] else 0,
                'Days': rep['num_days']
            })

    if not repeatability_data:
        return None, None

    df = pd.DataFrame(repeatability_data)
    df = df.sort_values('U (mm)')

    fig, ax = plt.subplots(figsize=(12, 6))
    x = np.arange(len(df))
    width = 0.25

    ax.bar(x - width, df['N (mm)'], width, label='North', color='#27ae60')
    ax.bar(x, df['E (mm)'], width, label='East', color='#3498db')
    ax.bar(x + width, df['U (mm)'], width, label='Up', color='#e74c3c')

    ax.set_xlabel('Station', fontsize=12)
    ax.set_ylabel('Repeatability (mm)', fontsize=12)
    ax.set_title(f'Coordinate Repeatability - DOY {start_doy}-{end_doy}/{year}', fontsize=14)
    ax.set_xticks(x)
    ax.set_xticklabels(df['Station'], rotation=45, ha='right', fontsize=8)
    ax.legend()
    ax.grid(axis='y', alpha=0.3)
    ax.axhline(y=3, color='orange', linestyle='--', alpha=0.7, label='3mm threshold')
    ax.axhline(y=5, color='red', linestyle='--', alpha=0.7, label='5mm threshold')

    plt.tight_layout()
    fig_path = FIG_DIR / "coord_repeatability.png"
    fig.savefig(fig_path, dpi=150, bbox_inches='tight')
    plt.close()

    return fig_path, df


def generate_rms_analysis_figure(db, year, start_doy, end_doy):
    """Generate RMS analysis figures"""
    stats = db.get_stats(year=year)
    df = pd.DataFrame(stats)

    if df.empty:
        return None, None

    df = df[(df['doy'] >= start_doy) & (df['doy'] <= end_doy)]
    df['rms_mm'] = df['rms_unit_weight'] * 1000

    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    # RMS time series
    ax1 = axes[0]
    for station in df['station_id'].unique()[:10]:
        station_df = df[df['station_id'] == station]
        ax1.plot(station_df['doy'], station_df['rms_mm'], 'o-', label=station, markersize=4)

    ax1.set_xlabel('DOY', fontsize=12)
    ax1.set_ylabel('RMS (mm)', fontsize=12)
    ax1.set_title('RMS Time Series by Station', fontsize=14)
    ax1.legend(fontsize=8, ncol=2)
    ax1.grid(alpha=0.3)
    ax1.axhline(y=2, color='green', linestyle='--', alpha=0.7)
    ax1.axhline(y=5, color='red', linestyle='--', alpha=0.7)

    # RMS histogram
    ax2 = axes[1]
    ax2.hist(df['rms_mm'], bins=30, color='#3498db', edgecolor='white', alpha=0.7)
    ax2.axvline(df['rms_mm'].mean(), color='red', linestyle='--', linewidth=2, label=f'Mean: {df["rms_mm"].mean():.2f} mm')
    ax2.axvline(df['rms_mm'].median(), color='green', linestyle='--', linewidth=2, label=f'Median: {df["rms_mm"].median():.2f} mm')
    ax2.set_xlabel('RMS (mm)', fontsize=12)
    ax2.set_ylabel('Frequency', fontsize=12)
    ax2.set_title('RMS Distribution', fontsize=14)
    ax2.legend()
    ax2.grid(alpha=0.3)

    plt.tight_layout()
    fig_path = FIG_DIR / "rms_analysis.png"
    fig.savefig(fig_path, dpi=150, bbox_inches='tight')
    plt.close()

    # Summary statistics
    summary = {
        'Metric': ['Mean RMS', 'Median RMS', 'Std RMS', 'Min RMS', 'Max RMS', '% < 2mm', '% < 5mm'],
        'Value': [
            f"{df['rms_mm'].mean():.2f} mm",
            f"{df['rms_mm'].median():.2f} mm",
            f"{df['rms_mm'].std():.2f} mm",
            f"{df['rms_mm'].min():.2f} mm",
            f"{df['rms_mm'].max():.2f} mm",
            f"{(df['rms_mm'] < 2).mean() * 100:.1f}%",
            f"{(df['rms_mm'] < 5).mean() * 100:.1f}%"
        ]
    }

    return fig_path, pd.DataFrame(summary)


def generate_ztd_analysis_figure(db, year, start_doy, end_doy):
    """Generate ZTD analysis figure"""
    # Get ZTD data for all stations
    conn = db.connect()
    result = conn.execute("""
        SELECT station_id, doy, hour, ztd, ztd_rms
        FROM ztd_hourly
        WHERE year = ? AND doy BETWEEN ? AND ?
        ORDER BY station_id, doy, hour
    """, [year, start_doy, end_doy]).fetchdf()

    if result.empty:
        return None, None

    df = result
    df['ztd_m'] = df['ztd']

    fig, axes = plt.subplots(2, 2, figsize=(14, 10))

    # ZTD time series for selected stations
    ax1 = axes[0, 0]
    stations_sample = df['station_id'].unique()[:5]
    for station in stations_sample:
        sdf = df[df['station_id'] == station]
        sdf['time'] = sdf['doy'] + sdf['hour']/24
        ax1.plot(sdf['time'], sdf['ztd_m'] * 1000, '-', label=station, linewidth=0.8)

    ax1.set_xlabel('DOY', fontsize=11)
    ax1.set_ylabel('ZTD (mm)', fontsize=11)
    ax1.set_title('ZTD Time Series (Sample Stations)', fontsize=12)
    ax1.legend(fontsize=8)
    ax1.grid(alpha=0.3)

    # Daily mean ZTD
    ax2 = axes[0, 1]
    daily_mean = df.groupby('doy')['ztd_m'].mean() * 1000
    ax2.bar(daily_mean.index, daily_mean.values, color='#3498db', alpha=0.7)
    ax2.set_xlabel('DOY', fontsize=11)
    ax2.set_ylabel('Mean ZTD (mm)', fontsize=11)
    ax2.set_title('Daily Mean ZTD (All Stations)', fontsize=12)
    ax2.grid(alpha=0.3)

    # ZTD distribution
    ax3 = axes[1, 0]
    ax3.hist(df['ztd_m'] * 1000, bins=50, color='#27ae60', edgecolor='white', alpha=0.7)
    ax3.set_xlabel('ZTD (mm)', fontsize=11)
    ax3.set_ylabel('Frequency', fontsize=11)
    ax3.set_title('ZTD Distribution', fontsize=12)
    ax3.grid(alpha=0.3)

    # Station mean ZTD
    ax4 = axes[1, 1]
    station_mean = df.groupby('station_id')['ztd_m'].mean().sort_values() * 1000
    ax4.barh(range(len(station_mean)), station_mean.values, color='#9b59b6', alpha=0.7)
    ax4.set_yticks(range(len(station_mean)))
    ax4.set_yticklabels(station_mean.index, fontsize=7)
    ax4.set_xlabel('Mean ZTD (mm)', fontsize=11)
    ax4.set_title('Station Mean ZTD', fontsize=12)
    ax4.grid(alpha=0.3)

    plt.tight_layout()
    fig_path = FIG_DIR / "ztd_analysis.png"
    fig.savefig(fig_path, dpi=150, bbox_inches='tight')
    plt.close()

    # Summary
    summary_df = pd.DataFrame({
        'Statistic': ['Total Entries', 'Stations', 'Mean ZTD', 'Std ZTD', 'Min ZTD', 'Max ZTD'],
        'Value': [
            len(df),
            df['station_id'].nunique(),
            f"{df['ztd_m'].mean()*1000:.1f} mm",
            f"{df['ztd_m'].std()*1000:.1f} mm",
            f"{df['ztd_m'].min()*1000:.1f} mm",
            f"{df['ztd_m'].max()*1000:.1f} mm"
        ]
    })

    return fig_path, summary_df


def generate_ambiguity_figure(db, year, start_doy, end_doy):
    """Generate ambiguity resolution analysis"""
    amb_data = db.get_ambiguity(year=year, start_doy=start_doy, end_doy=end_doy)

    if not amb_data:
        return None, None

    df = pd.DataFrame(amb_data)

    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    # WL/NL rates by station
    ax1 = axes[0]
    stations = df['station_id'].unique()[:15]
    station_df = df[df['station_id'].isin(stations)]
    wl_means = station_df.groupby('station_id')['wl_combined'].mean()
    nl_means = station_df.groupby('station_id')['nl_combined'].mean()

    x = np.arange(len(wl_means))
    width = 0.35
    ax1.bar(x - width/2, wl_means.values, width, label='Widelane', color='#3498db')
    ax1.bar(x + width/2, nl_means.values, width, label='Narrowlane', color='#e74c3c')
    ax1.set_xticks(x)
    ax1.set_xticklabels(wl_means.index, rotation=45, ha='right', fontsize=8)
    ax1.set_ylabel('Resolution Rate (%)', fontsize=11)
    ax1.set_title('Ambiguity Resolution by Station', fontsize=12)
    ax1.legend()
    ax1.grid(alpha=0.3)
    ax1.axhline(y=80, color='green', linestyle='--', alpha=0.7)

    # Time series
    ax2 = axes[1]
    daily_wl = df.groupby('doy')['wl_combined'].mean()
    daily_nl = df.groupby('doy')['nl_combined'].mean()
    ax2.plot(daily_wl.index, daily_wl.values, 'o-', label='Widelane', color='#3498db', markersize=6)
    ax2.plot(daily_nl.index, daily_nl.values, 's-', label='Narrowlane', color='#e74c3c', markersize=6)
    ax2.set_xlabel('DOY', fontsize=11)
    ax2.set_ylabel('Resolution Rate (%)', fontsize=11)
    ax2.set_title('Daily Ambiguity Resolution Rates', fontsize=12)
    ax2.legend()
    ax2.grid(alpha=0.3)
    ax2.axhline(y=80, color='green', linestyle='--', alpha=0.7)

    plt.tight_layout()
    fig_path = FIG_DIR / "ambiguity_resolution.png"
    fig.savefig(fig_path, dpi=150, bbox_inches='tight')
    plt.close()

    # Summary
    summary_df = pd.DataFrame({
        'Metric': ['Mean WL Rate', 'Mean NL Rate', 'Mean GPS WL', 'Mean GAL WL', 'Stations', 'Days'],
        'Value': [
            f"{df['wl_combined'].mean():.1f}%",
            f"{df['nl_combined'].mean():.1f}%",
            f"{df['wl_gps'].mean():.1f}%" if 'wl_gps' in df else "N/A",
            f"{df['wl_gal'].mean():.1f}%" if 'wl_gal' in df else "N/A",
            df['station_id'].nunique(),
            df['doy'].nunique()
        ]
    })

    return fig_path, summary_df


def generate_satellite_tracking_figure(db, year, start_doy, end_doy):
    """Generate satellite tracking analysis"""
    sat_data = db.get_satellite_tracking(year=year)

    if not sat_data:
        return None, None

    df = pd.DataFrame(sat_data)
    df = df[(df['doy'] >= start_doy) & (df['doy'] <= end_doy)]

    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    # RMS by constellation
    ax1 = axes[0]
    for const in df['constellation'].unique():
        const_df = df[df['constellation'] == const]
        prn_rms = const_df.groupby('prn')['rms'].mean()
        color = {'GPS': '#27ae60', 'GALILEO': '#3498db', 'GLONASS': '#e74c3c'}.get(const, '#95a5a6')
        ax1.bar(prn_rms.index + (list(df['constellation'].unique()).index(const) * 0.25),
                prn_rms.values, width=0.25, label=const, color=color, alpha=0.7)

    ax1.set_xlabel('PRN', fontsize=11)
    ax1.set_ylabel('RMS (mm)', fontsize=11)
    ax1.set_title('Satellite RMS by PRN', fontsize=12)
    ax1.legend()
    ax1.grid(alpha=0.3)

    # Observation percentage
    ax2 = axes[1]
    const_obs = df.groupby('constellation')['obs_percent'].mean()
    colors = [{'GPS': '#27ae60', 'GALILEO': '#3498db', 'GLONASS': '#e74c3c'}.get(c, '#95a5a6') for c in const_obs.index]
    ax2.bar(const_obs.index, const_obs.values, color=colors, alpha=0.7)
    ax2.set_ylabel('Mean Observation %', fontsize=11)
    ax2.set_title('Mean Observation % by Constellation', fontsize=12)
    ax2.grid(alpha=0.3)

    plt.tight_layout()
    fig_path = FIG_DIR / "satellite_tracking.png"
    fig.savefig(fig_path, dpi=150, bbox_inches='tight')
    plt.close()

    # Summary by constellation
    summary_df = df.groupby('constellation').agg({
        'prn': 'nunique',
        'rms': 'mean',
        'obs_percent': 'mean',
        'obs_count': 'sum'
    }).reset_index()
    summary_df.columns = ['Constellation', 'Satellites', 'Mean RMS (mm)', 'Mean Obs %', 'Total Obs']
    summary_df['Mean RMS (mm)'] = summary_df['Mean RMS (mm)'].round(2)
    summary_df['Mean Obs %'] = summary_df['Mean Obs %'].round(1)

    return fig_path, summary_df


def generate_code_coordinate_comparison(db, year, start_doy, end_doy):
    """Generate CODE coordinate comparison"""
    stations = db.get_all_stations()
    comparison_data = []

    for doy in range(start_doy, end_doy + 1):
        # Try to download/load CODE SNX
        snx_path = download_code_snx(year, doy, product_type="FIN")
        if not snx_path:
            continue

        try:
            code_coords = parse_code_snx_file(snx_path)
        except:
            continue

        # Get local solutions
        for station in stations:
            sols = db.get_solutions(station_id=station, year=year, start_doy=doy, end_doy=doy)
            if not sols:
                continue

            sol = sols[0]
            code_coord = get_code_coords_for_station(code_coords, station)
            if not code_coord:
                continue

            # Compute ENU differences
            dE, dN, dU = compute_enu_diff(sol['x'], sol['y'], sol['z'], code_coord)

            comparison_data.append({
                'Station': station,
                'DOY': doy,
                'dE (mm)': dE,
                'dN (mm)': dN,
                'dU (mm)': dU,
                '3D (mm)': np.sqrt(dE**2 + dN**2 + dU**2)
            })

    if not comparison_data:
        return None, None

    df = pd.DataFrame(comparison_data)

    fig, axes = plt.subplots(2, 2, figsize=(14, 10))

    # ENU by station
    ax1 = axes[0, 0]
    station_means = df.groupby('Station')[['dE (mm)', 'dN (mm)', 'dU (mm)']].mean()
    x = np.arange(len(station_means))
    width = 0.25
    ax1.bar(x - width, station_means['dE (mm)'], width, label='East', color='#3498db')
    ax1.bar(x, station_means['dN (mm)'], width, label='North', color='#27ae60')
    ax1.bar(x + width, station_means['dU (mm)'], width, label='Up', color='#e74c3c')
    ax1.set_xticks(x)
    ax1.set_xticklabels(station_means.index, rotation=45, ha='right', fontsize=7)
    ax1.set_ylabel('Difference (mm)', fontsize=11)
    ax1.set_title('Mean ENU Differences vs CODE (by Station)', fontsize=12)
    ax1.legend()
    ax1.grid(alpha=0.3)
    ax1.axhline(y=0, color='black', linestyle='-', linewidth=0.5)

    # Time series
    ax2 = axes[0, 1]
    daily_mean = df.groupby('DOY')[['dE (mm)', 'dN (mm)', 'dU (mm)']].mean()
    ax2.plot(daily_mean.index, daily_mean['dE (mm)'], 'o-', label='East', color='#3498db')
    ax2.plot(daily_mean.index, daily_mean['dN (mm)'], 's-', label='North', color='#27ae60')
    ax2.plot(daily_mean.index, daily_mean['dU (mm)'], '^-', label='Up', color='#e74c3c')
    ax2.set_xlabel('DOY', fontsize=11)
    ax2.set_ylabel('Mean Difference (mm)', fontsize=11)
    ax2.set_title('Daily Mean ENU Differences vs CODE', fontsize=12)
    ax2.legend()
    ax2.grid(alpha=0.3)
    ax2.axhline(y=0, color='black', linestyle='-', linewidth=0.5)

    # 3D RMS histogram
    ax3 = axes[1, 0]
    ax3.hist(df['3D (mm)'], bins=30, color='#9b59b6', edgecolor='white', alpha=0.7)
    ax3.axvline(df['3D (mm)'].mean(), color='red', linestyle='--', linewidth=2, label=f'Mean: {df["3D (mm)"].mean():.1f} mm')
    ax3.set_xlabel('3D Difference (mm)', fontsize=11)
    ax3.set_ylabel('Frequency', fontsize=11)
    ax3.set_title('3D Position Difference Distribution', fontsize=12)
    ax3.legend()
    ax3.grid(alpha=0.3)

    # RMS by station
    ax4 = axes[1, 1]
    station_rms = df.groupby('Station')['3D (mm)'].apply(lambda x: np.sqrt((x**2).mean())).sort_values()
    ax4.barh(range(len(station_rms)), station_rms.values, color='#f39c12', alpha=0.7)
    ax4.set_yticks(range(len(station_rms)))
    ax4.set_yticklabels(station_rms.index, fontsize=7)
    ax4.set_xlabel('3D RMS (mm)', fontsize=11)
    ax4.set_title('3D RMS vs CODE by Station', fontsize=12)
    ax4.grid(alpha=0.3)
    ax4.axvline(x=5, color='green', linestyle='--', alpha=0.7)
    ax4.axvline(x=10, color='red', linestyle='--', alpha=0.7)

    plt.tight_layout()
    fig_path = FIG_DIR / "code_coordinate_comparison.png"
    fig.savefig(fig_path, dpi=150, bbox_inches='tight')
    plt.close()

    # Summary statistics
    summary_df = pd.DataFrame({
        'Component': ['East', 'North', 'Up', '3D'],
        'Mean (mm)': [df['dE (mm)'].mean(), df['dN (mm)'].mean(), df['dU (mm)'].mean(), df['3D (mm)'].mean()],
        'Std (mm)': [df['dE (mm)'].std(), df['dN (mm)'].std(), df['dU (mm)'].std(), df['3D (mm)'].std()],
        'RMS (mm)': [
            np.sqrt((df['dE (mm)']**2).mean()),
            np.sqrt((df['dN (mm)']**2).mean()),
            np.sqrt((df['dU (mm)']**2).mean()),
            np.sqrt((df['3D (mm)']**2).mean())
        ]
    })
    summary_df = summary_df.round(2)

    return fig_path, summary_df


def generate_code_ztd_comparison(db, year, start_doy, end_doy):
    """Generate CODE ZTD comparison"""
    comparison_data = []

    for doy in range(start_doy, end_doy + 1):
        # Try to download/load CODE TRO
        tro_path = download_code_tro(year, doy, product_type="FIN")
        if not tro_path:
            continue

        try:
            _, code_entries = parse_code_tro_file(tro_path)
        except:
            continue

        # Get local ZTD
        conn = db.connect()
        local_ztd = conn.execute("""
            SELECT station_id, hour, ztd, ztd_rms
            FROM ztd_hourly
            WHERE year = ? AND doy = ?
        """, [year, doy]).fetchdf()

        if local_ztd.empty:
            continue

        for _, row in local_ztd.iterrows():
            code_entries_station = get_code_ztd_for_station(code_entries, row['station_id'], year, doy)
            # Find matching hour
            for entry in code_entries_station:
                if entry.hour == row['hour']:
                    diff = (row['ztd'] - entry.ztd) * 1000  # mm
                    comparison_data.append({
                        'Station': row['station_id'],
                        'DOY': doy,
                        'Hour': row['hour'],
                        'Local ZTD (mm)': row['ztd'] * 1000,
                        'CODE ZTD (mm)': entry.ztd * 1000,
                        'Diff (mm)': diff
                    })
                    break

    if not comparison_data:
        return None, None

    df = pd.DataFrame(comparison_data)

    fig, axes = plt.subplots(2, 2, figsize=(14, 10))

    # ZTD difference by station
    ax1 = axes[0, 0]
    station_means = df.groupby('Station')['Diff (mm)'].mean().sort_values()
    colors = ['#27ae60' if abs(v) < 3 else '#f39c12' if abs(v) < 5 else '#e74c3c' for v in station_means.values]
    ax1.barh(range(len(station_means)), station_means.values, color=colors, alpha=0.7)
    ax1.set_yticks(range(len(station_means)))
    ax1.set_yticklabels(station_means.index, fontsize=7)
    ax1.set_xlabel('Mean ZTD Difference (mm)', fontsize=11)
    ax1.set_title('ZTD Difference vs CODE by Station', fontsize=12)
    ax1.grid(alpha=0.3)
    ax1.axvline(x=0, color='black', linestyle='-', linewidth=0.5)

    # Time series
    ax2 = axes[0, 1]
    daily_mean = df.groupby('DOY')['Diff (mm)'].mean()
    daily_std = df.groupby('DOY')['Diff (mm)'].std()
    ax2.errorbar(daily_mean.index, daily_mean.values, yerr=daily_std.values,
                 fmt='o-', color='#3498db', capsize=3, markersize=6)
    ax2.set_xlabel('DOY', fontsize=11)
    ax2.set_ylabel('Mean ZTD Difference (mm)', fontsize=11)
    ax2.set_title('Daily Mean ZTD Difference vs CODE', fontsize=12)
    ax2.grid(alpha=0.3)
    ax2.axhline(y=0, color='black', linestyle='-', linewidth=0.5)

    # Histogram
    ax3 = axes[1, 0]
    ax3.hist(df['Diff (mm)'], bins=50, color='#9b59b6', edgecolor='white', alpha=0.7)
    ax3.axvline(df['Diff (mm)'].mean(), color='red', linestyle='--', linewidth=2,
                label=f'Mean: {df["Diff (mm)"].mean():.2f} mm')
    ax3.set_xlabel('ZTD Difference (mm)', fontsize=11)
    ax3.set_ylabel('Frequency', fontsize=11)
    ax3.set_title('ZTD Difference Distribution', fontsize=12)
    ax3.legend()
    ax3.grid(alpha=0.3)

    # Scatter plot
    ax4 = axes[1, 1]
    ax4.scatter(df['CODE ZTD (mm)'], df['Local ZTD (mm)'], alpha=0.3, s=10, c='#3498db')
    lims = [min(df['CODE ZTD (mm)'].min(), df['Local ZTD (mm)'].min()),
            max(df['CODE ZTD (mm)'].max(), df['Local ZTD (mm)'].max())]
    ax4.plot(lims, lims, 'r--', linewidth=1, label='1:1 line')
    ax4.set_xlabel('CODE ZTD (mm)', fontsize=11)
    ax4.set_ylabel('Local ZTD (mm)', fontsize=11)
    ax4.set_title('Local vs CODE ZTD Comparison', fontsize=12)
    ax4.legend()
    ax4.grid(alpha=0.3)

    plt.tight_layout()
    fig_path = FIG_DIR / "code_ztd_comparison.png"
    fig.savefig(fig_path, dpi=150, bbox_inches='tight')
    plt.close()

    # Summary
    summary_df = pd.DataFrame({
        'Metric': ['Comparisons', 'Stations', 'Days', 'Mean Diff', 'Std Diff', 'RMS Diff', '% < 3mm', '% < 5mm'],
        'Value': [
            len(df),
            df['Station'].nunique(),
            df['DOY'].nunique(),
            f"{df['Diff (mm)'].mean():.2f} mm",
            f"{df['Diff (mm)'].std():.2f} mm",
            f"{np.sqrt((df['Diff (mm)']**2).mean()):.2f} mm",
            f"{(abs(df['Diff (mm)']) < 3).mean()*100:.1f}%",
            f"{(abs(df['Diff (mm)']) < 5).mean()*100:.1f}%"
        ]
    })

    return fig_path, summary_df


def generate_presentation(year=2025, start_doy=296, end_doy=310, output_path=None):
    """Generate the complete PowerPoint presentation"""

    if output_path is None:
        output_path = f"/home/ahunegnaw/GNSS_Quality_Report_{year}_DOY{start_doy}-{end_doy}.pptx"

    print(f"Generating GNSS Quality Report for {year} DOY {start_doy}-{end_doy}")

    # Initialize database
    db = QualityDatabase(DB_PATH)
    db.connect()

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== TITLE SLIDE ==========
    add_title_slide(
        prs,
        "GNSS Quality Monitoring Report",
        f"Year {year} | DOY {start_doy}-{end_doy} | Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )

    # ========== EXECUTIVE SUMMARY ==========
    print("Generating executive summary...")
    stations = db.get_all_stations()
    stats = db.get_stats(year=year)
    stats_df = pd.DataFrame(stats)
    stats_df = stats_df[(stats_df['doy'] >= start_doy) & (stats_df['doy'] <= end_doy)]

    summary_items = [
        f"Processing Period: DOY {start_doy} - {end_doy}, Year {year}",
        f"Total Stations: {len(stations)}",
        f"Total Processing Records: {len(stats_df)}",
        f"Days Processed: {stats_df['doy'].nunique()}",
        "",
        "Key Findings:",
        f"  Mean RMS: {stats_df['rms_unit_weight'].mean()*1000:.2f} mm",
        f"  Solutions with RMS < 2mm: {(stats_df['rms_unit_weight']*1000 < 2).mean()*100:.1f}%",
        f"  Total Observations: {stats_df['num_observations'].sum():,}",
        "",
        "Report Contents:",
        "  1. Coordinate Repeatability Analysis",
        "  2. RMS & Processing Statistics",
        "  3. Troposphere (ZTD) Analysis",
        "  4. Ambiguity Resolution Statistics",
        "  5. Satellite Tracking Quality",
        "  6. CODE Product Comparisons"
    ]
    add_content_slide(prs, "Executive Summary", summary_items)

    # ========== COORDINATE REPEATABILITY ==========
    print("Generating coordinate repeatability...")
    add_section_slide(prs, "1. Coordinate Repeatability")

    fig_path, rep_df = generate_coordinate_repeatability_figure(db, year, start_doy, end_doy)
    if fig_path:
        add_content_slide(prs, "Coordinate Repeatability (NEU)", image_path=fig_path)

        # Stats table
        if rep_df is not None and not rep_df.empty:
            rep_summary = rep_df.describe()[['N (mm)', 'E (mm)', 'U (mm)']].round(2)
            add_table_slide(prs, "Coordinate Repeatability Statistics", rep_summary.reset_index())

    # ========== RMS ANALYSIS ==========
    print("Generating RMS analysis...")
    add_section_slide(prs, "2. RMS Analysis")

    fig_path, rms_summary = generate_rms_analysis_figure(db, year, start_doy, end_doy)
    if fig_path:
        add_content_slide(prs, "RMS Time Series & Distribution", image_path=fig_path)
        if rms_summary is not None:
            add_table_slide(prs, "RMS Summary Statistics", rms_summary)

    # ========== ZTD ANALYSIS ==========
    print("Generating ZTD analysis...")
    add_section_slide(prs, "3. Troposphere (ZTD) Analysis")

    fig_path, ztd_summary = generate_ztd_analysis_figure(db, year, start_doy, end_doy)
    if fig_path:
        add_content_slide(prs, "ZTD Time Series & Distribution", image_path=fig_path)
        if ztd_summary is not None:
            add_table_slide(prs, "ZTD Summary Statistics", ztd_summary)

    # ========== AMBIGUITY RESOLUTION ==========
    print("Generating ambiguity resolution analysis...")
    add_section_slide(prs, "4. Ambiguity Resolution")

    fig_path, amb_summary = generate_ambiguity_figure(db, year, start_doy, end_doy)
    if fig_path:
        add_content_slide(prs, "Ambiguity Resolution Rates", image_path=fig_path)
        if amb_summary is not None:
            add_table_slide(prs, "Ambiguity Resolution Summary", amb_summary)

    # ========== SATELLITE TRACKING ==========
    print("Generating satellite tracking analysis...")
    add_section_slide(prs, "5. Satellite Tracking")

    fig_path, sat_summary = generate_satellite_tracking_figure(db, year, start_doy, end_doy)
    if fig_path:
        add_content_slide(prs, "Satellite Tracking Quality", image_path=fig_path)
        if sat_summary is not None:
            add_table_slide(prs, "Satellite Tracking by Constellation", sat_summary)

    # ========== CODE COMPARISONS ==========
    print("Generating CODE coordinate comparison...")
    add_section_slide(prs, "6. Comparison with CODE Products")

    # Coordinate comparison
    fig_path, coord_summary = generate_code_coordinate_comparison(db, year, start_doy, end_doy)
    if fig_path:
        add_content_slide(prs, "Coordinate Comparison with CODE SINEX", image_path=fig_path)
        if coord_summary is not None:
            add_table_slide(prs, "ENU Differences vs CODE (Statistics)", coord_summary)
    else:
        add_content_slide(prs, "Coordinate Comparison with CODE SINEX",
                         ["No CODE coordinate products available for comparison",
                          "Products are downloaded from ftp.aiub.unibe.ch/CODE/"])

    # ZTD comparison
    print("Generating CODE ZTD comparison...")
    fig_path, ztd_code_summary = generate_code_ztd_comparison(db, year, start_doy, end_doy)
    if fig_path:
        add_content_slide(prs, "ZTD Comparison with CODE Troposphere", image_path=fig_path)
        if ztd_code_summary is not None:
            add_table_slide(prs, "ZTD Differences vs CODE (Statistics)", ztd_code_summary)
    else:
        add_content_slide(prs, "ZTD Comparison with CODE Troposphere",
                         ["No CODE troposphere products available for comparison",
                          "Products are downloaded from ftp.aiub.unibe.ch/CODE/"])

    # ========== CONCLUSIONS ==========
    add_section_slide(prs, "Conclusions & Recommendations")

    conclusions = [
        "Processing Quality Assessment:",
        f"  Overall RMS quality: {'Good' if stats_df['rms_unit_weight'].mean()*1000 < 3 else 'Moderate' if stats_df['rms_unit_weight'].mean()*1000 < 5 else 'Needs Review'}",
        f"  {len(stations)} stations processed successfully over {end_doy-start_doy+1} days",
        "",
        "Recommendations:",
        "  Review stations with coordinate repeatability > 5mm",
        "  Investigate days with anomalous RMS values",
        "  Monitor ambiguity resolution rates below 80%",
        "",
        "Data Sources:",
        "  Local processing: Bernese GNSS Software PPP",
        "  Reference: CODE final products (ftp.aiub.unibe.ch)",
        "",
        f"Report generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    ]
    add_content_slide(prs, "Conclusions & Recommendations", conclusions)

    # Save presentation
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")

    return output_path


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Generate GNSS Quality Report PowerPoint")
    parser.add_argument("--year", type=int, default=2025, help="Year")
    parser.add_argument("--start-doy", type=int, default=296, help="Start DOY")
    parser.add_argument("--end-doy", type=int, default=310, help="End DOY")
    parser.add_argument("--output", type=str, default=None, help="Output file path")

    args = parser.parse_args()

    generate_presentation(args.year, args.start_doy, args.end_doy, args.output)
