#!/usr/bin/env python3
"""
Create PowerPoint presentation for Advanced PPP-AR Strategies
using Bernese GNSS Software v5.4

Author: Addisu Hunegnaw
Date: January 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, sub_bullets=None):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 102, 153)
    line.line.fill.background()

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        p.font.color.rgb = RGBColor(50, 50, 50)

        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                p = tf.add_paragraph()
                p.text = f"    ◦ {sub}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(80, 80, 80)
                p.space_after = Pt(6)

    return slide


def add_table_slide(prs, title, headers, rows, notes=None):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1

    x, y, cx, cy = Inches(0.5), Inches(1.3), Inches(9), Inches(0.4 * table_rows)
    table = slide.shapes.add_table(table_rows, cols, x, y, cx, cy).table

    # Set column widths
    col_width = Inches(9 / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)

    # Data rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 240, 250)

    # Notes below table
    if notes:
        notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5 + 0.4 * table_rows), Inches(9), Inches(2))
        tf = notes_box.text_frame
        tf.word_wrap = True
        for i, note in enumerate(notes):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = note
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(80, 80, 80)

    return slide


def add_code_slide(prs, title, code_lines, description=None):
    """Add a slide with code/config display"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Description
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.6))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(80, 80, 80)
        code_y = 1.8
    else:
        code_y = 1.3

    # Code box (with background)
    code_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(code_y), Inches(9), Inches(5.5 - code_y + 1))
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = RGBColor(40, 44, 52)
    code_shape.line.fill.background()

    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(code_y + 0.2), Inches(8.6), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(171, 178, 191)
        p.space_after = Pt(2)

    return slide


def create_presentation():
    """Create the full PPP-AR presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "Advanced PPP-AR Strategies using\nBernese GNSS Software v5.4",
        "High-Precision Positioning with Integer Ambiguity Resolution\n\nAddisu Hunegnaw | January 2026"
    )

    # ========== SLIDE 2: Overview ==========
    add_content_slide(prs, "Presentation Overview", [
        "Processing Strategy & BPE Configuration",
        "Multi-GNSS Observation Handling (GPS, Galileo, GLONASS, BeiDou)",
        "High-Rate Clock and Phase Bias Products",
        "Zero-Difference Ambiguity Resolution Techniques",
        "Critical Modeling Components (APC, OTL, Troposphere)",
        "Quality Control: Ocean Loading & Satellite Residuals",
        "GLONASS FDMA Limitations"
    ])

    # ========== SLIDE 3: PPP-AR Fundamentals ==========
    add_content_slide(prs, "PPP-AR Fundamentals", [
        "PPP (Precise Point Positioning): Single-receiver positioning using precise orbit/clock products",
        "AR (Ambiguity Resolution): Fixing carrier phase ambiguities to integer values",
        "Benefits of PPP-AR:",
        "Carrier Phase Observation: L = ρ + c(dt_r - dt_s) + T + I + λN + ε"
    ], sub_bullets={
        2: ["Convergence time: 30+ min → 5-10 min",
            "Horizontal accuracy: 2-3 cm → 5-10 mm",
            "Vertical accuracy: 5-6 cm → 15-20 mm"]
    })

    # ========== SLIDE 4: BPE Workflow ==========
    add_code_slide(prs, "Bernese 5.4 Processing Chain Overview",
        [
            "PID  SCRIPT    DESCRIPTION",
            "─────────────────────────────────────────────",
            "001  EGV_COP   → Copy required files",
            "101  POLUPD    → Update pole information",
            "111  ORBMRG    → Merge orbit files",
            "211  RNXSMT    → RINEX smoothing",
            "221  RXOBV3    → Observation file creation",
            "231  CODSPP    → Code-based SPP",
            "301-342        → Baseline formation & preprocessing",
            "401-443        → Ambiguity resolution (MW, L53, QIF, L12)",
            "501-531        → Final solution & quality assessment",
            "901  R2S_SUM   → Summary and coordinate repeatability"
        ],
        "BPE (Bernese Processing Engine) executes programs in sequence defined by PCF"
    )

    # ========== SLIDE 5: Multi-GNSS Configuration ==========
    add_table_slide(prs, "Multi-GNSS Configuration",
        ["Variable", "Value", "Description"],
        [
            ["V_SATSYS", "GRE", "Processing: GPS + GLONASS + Galileo"],
            ["V_GNSSAR", "GRE", "Ambiguity resolution systems"],
            ["V_PCV", "I20", "IGS20 antenna phase center model"],
            ["V_ORB", "COD", "CODE orbit/clock products"],
            ["V_SAMPL", "180", "3-minute sampling interval"]
        ],
        ["Critical Distinction:",
         "• V_SATSYS: Which satellites to process",
         "• V_GNSSAR: Which systems for AR (GLONASS excluded in practice due to FDMA)"]
    )

    # ========== SLIDE 6: Orbit and Clock Products ==========
    add_table_slide(prs, "CODE Analysis Center Products",
        ["Product", "Latency", "Orbit Accuracy", "Clock Accuracy"],
        [
            ["COD Final", "12-18 days", "2.5 cm", "75 ps"],
            ["COD Rapid", "17-41 hours", "2.5 cm", "75 ps"],
            ["COD Ultra-Rapid", "3-9 hours", "5 cm", "150 ps"]
        ],
        ["High-Rate Clocks for PPP-AR:",
         "• 30-second clock sampling essential for PPP",
         "• Phase-consistent clocks enable ambiguity resolution",
         "• OSB files (IAR_wwwwd.OSB) provide satellite phase biases"]
    )

    # ========== SLIDE 7: Observable-Specific Biases ==========
    add_content_slide(prs, "Observable-Specific Biases (OSB)", [
        "OSB Files: IAR_wwwwd.OSB from CODE analysis center",
        "Purpose: Enable integer ambiguity recovery in undifferenced PPP",
        "Bias Application: φ_corrected = φ_raw - OSB_satellite - OSB_receiver",
        "Key Signals Supported:"
    ], sub_bullets={
        3: ["GPS: L1C, L2W, L5Q",
            "Galileo: E1C, E5a, E5b, E6C",
            "BeiDou: B1I, B2I, B3I"]
    })

    # ========== SLIDE 8: Ambiguity Resolution Strategy ==========
    add_content_slide(prs, "Four-Step Ambiguity Resolution", [
        "Step 1 - Melbourne-Wübbena (MW): Wide-lane λ ≈ 86 cm, BL limit: 6000 km",
        "Step 2 - L5/L3 Resolution: Uses resolved WL to constrain L3, BL limit: 200 km",
        "Step 3 - QIF (Quasi-Ionosphere-Free): Direct L1/L2 resolution, BL limit: 2000 km",
        "Step 4 - L1 & L2 Direct: Final integer resolution, BL limit: 20 km",
        "Cascading approach: Each step uses constraints from previous steps"
    ])

    # ========== SLIDE 9: GLONASS FDMA Limitation ==========
    add_table_slide(prs, "GLONASS FDMA Limitation",
        ["System", "Modulation", "Frequency", "AR Status"],
        [
            ["GPS", "CDMA", "L1: 1575.42 MHz (all sats)", "✓ Supported"],
            ["Galileo", "CDMA", "E1: 1575.42 MHz (all sats)", "✓ Supported"],
            ["BeiDou", "CDMA", "B1: 1561.098 MHz (all sats)", "✓ Supported"],
            ["GLONASS", "FDMA", "L1: 1602 + k×0.5625 MHz", "✗ Not Supported"]
        ],
        ["FDMA Impact:",
         "• 14 different frequencies (k = -7 to +6)",
         "• Inter-Channel Biases (ICBs) vary by receiver type",
         "• ICBs cannot be reliably calibrated for AR",
         "• GLONASS still improves geometry in float solution"]
    )

    # ========== SLIDE 10: Antenna Phase Center Modeling ==========
    add_content_slide(prs, "Antenna Phase Center Modeling (IGS20)", [
        "Type-Mean Calibrations: Average over antenna population",
        "Individual Calibrations: Robot/chamber calibrated (preferred)",
        "Correction Components:",
        "Multi-GNSS Calibrations: Separate values for each frequency"
    ], sub_bullets={
        2: ["PCO: Phase Center Offset (North, East, Up)",
            "PCV: Phase Center Variations (elevation/azimuth dependent)"],
        3: ["GPS L1/L2/L5", "Galileo E1/E5a/E5b/E6", "GLONASS G1/G2"]
    })

    # ========== SLIDE 11: Tropospheric Modeling ==========
    add_content_slide(prs, "Tropospheric Modeling (VMF3)", [
        "Hydrostatic delay: ~2.3 m at zenith (well-modeled)",
        "Wet delay: 0-40 cm (highly variable, must be estimated)",
        "VMF3 Mapping Functions: Elevation-dependent + azimuthal asymmetry",
        "Estimation Strategy:",
        "Products: ZTD estimates can be extracted for meteorological applications"
    ], sub_bullets={
        3: ["A priori: VMF3 + GPT3 model (90%+ of total delay)",
            "Estimate: ZTD every 1-2 hours",
            "Optional: Horizontal gradients (North, East)"]
    })

    # ========== SLIDE 12: Ocean Tide Loading ==========
    add_content_slide(prs, "Ocean Tide Loading (FES2014b)", [
        "Effect: Crustal deformation from gravitational attraction of ocean mass",
        "Magnitude: Up to 10+ cm vertical displacement in coastal regions",
        "11 Tidal Constituents modeled:",
        "BLQ File: Contains amplitude and phase for each constituent per station"
    ], sub_bullets={
        2: ["Semi-diurnal: M2, S2, N2, K2 (periods ~12 hours)",
            "Diurnal: K1, O1, P1, Q1 (periods ~24 hours)",
            "Long-period: Mf, Mm, Ssa (fortnightly to semi-annual)"]
    })

    # ========== SLIDE 13: QC - Ocean Loading Analysis for PPP-AR ==========
    add_content_slide(prs, "QC: Ocean Loading Analysis for PPP-AR", [
        "Purpose: Verify OTL corrections are properly applied in solutions",
        "Dominant Constituent (M2): Principal lunar semi-diurnal, period 12.42 hours",
        "Maximum Displacements by Region:",
        "Validation Approach:",
        "Residual Signature: Unmodeled OTL shows ~12-hour periodicity in position residuals"
    ], sub_bullets={
        2: ["Coastal Europe (BRST, NEWL): Up to 50+ mm vertical",
            "Inland Europe (POTS, WTZR): 10-20 mm vertical",
            "Island sites (REYK): Strong horizontal components"],
        3: ["Compare estimated ZTD with unloaded positions",
            "Check for systematic periodic residuals at tidal frequencies",
            "Verify BLQ coefficients from Onsala loading service"]
    })

    # ========== SLIDE 14: QC - Satellite Residuals for PPP-AR ==========
    add_table_slide(prs, "QC: Satellite Residuals for PPP-AR",
        ["Metric", "Target", "Warning", "Critical"],
        [
            ["Phase RMS (L3)", "< 3 mm", "3-5 mm", "> 5 mm"],
            ["Code RMS (P3)", "< 50 cm", "50-80 cm", "> 100 cm"],
            ["Single Satellite RMS", "< 2 mm", "2-3 mm", "> 3 mm"],
            ["Observation Count", "> 95%", "90-95%", "< 90%"]
        ],
        ["Satellite-Specific Analysis:",
         "• Identify consistently high-residual satellites (possible orbit/clock issues)",
         "• Compare GPS vs Galileo vs GLONASS performance",
         "• Flag satellites with RMS > 1.5× network mean",
         "• Monitor residuals from RES_*.SUM.gz files per session"]
    )

    # ========== SLIDE 15: QC Metrics Summary ==========
    add_table_slide(prs, "Quality Control Metrics Summary",
        ["Metric", "Target Value", "Critical Threshold"],
        [
            ["Ambiguity Fixing Rate", "> 90%", "< 70%"],
            ["WRMS North", "< 2 mm", "> 5 mm"],
            ["WRMS East", "< 2 mm", "> 5 mm"],
            ["WRMS Up", "< 5 mm", "> 10 mm"],
            ["Convergence Time (AR)", "< 15 min", "> 30 min"],
            ["Phase Residual RMS", "< 3 mm", "> 5 mm"]
        ],
        ["Monitoring Files: AMB_*.SUM (ambiguity), HLM_*.OUT (Helmert),",
         "R2S_*.PRC (repeatability), RES_*.SUM (residuals), TRP files (convergence)"]
    )

    # ========== SLIDE 16: PPP-AR Coordinate Repeatability (Part 1) ==========
    add_table_slide(prs, "PPP-AR Repeatability - All 38 Stations (1/2)",
        ["Station", "Days", "E(mm)", "N(mm)", "U(mm)", "", "Station", "Days", "E(mm)", "N(mm)", "U(mm)"],
        [
            ["0531", "15", "2.33", "2.19", "7.47", "", "LROC", "15", "1.94", "2.04", "5.41"],
            ["0929", "15", "2.51", "1.68", "9.27", "", "MABO", "15", "1.84", "1.67", "7.60"],
            ["0935", "15", "2.10", "1.53", "6.97", "", "MAR6", "15", "3.06", "1.75", "5.80"],
            ["ARL0", "15", "2.85", "1.46", "4.85", "", "MATE", "15", "0.72", "1.43", "4.35"],
            ["BASC", "15", "2.56", "1.95", "8.19", "", "MLVL", "15", "1.99", "1.57", "5.75"],
            ["BOR1", "15", "1.73", "1.92", "9.89", "", "ONSA", "15", "1.79", "2.16", "11.46"],
            ["BRST", "15", "1.67", "2.35", "5.63", "", "POTS", "15", "1.60", "1.70", "10.16"],
            ["BRUX", "15", "2.46", "1.38", "8.31", "", "ROUL", "15", "2.95", "2.29", "7.30"],
            ["DARE", "15", "2.91", "1.25", "6.56", "", "STAS", "15", "2.13", "2.30", "10.73"],
            ["ECH2", "15", "2.59", "2.58", "7.41", "", "THNV", "15", "2.30", "1.11", "3.12"],
        ],
        ["Source: CRD files, Year 2025, DOY 296-310 (reprocessed 24-Jan-2026)"]
    )

    # ========== SLIDE 17: PPP-AR Coordinate Repeatability (Part 2) ==========
    add_table_slide(prs, "PPP-AR Repeatability - All 38 Stations (2/2)",
        ["Station", "Days", "E(mm)", "N(mm)", "U(mm)", "", "Station", "Days", "E(mm)", "N(mm)", "U(mm)"],
        [
            ["ENTZ", "15", "1.40", "1.89", "9.85", "", "TLSE", "15", "1.50", "1.31", "4.72"],
            ["ERPL", "15", "2.45", "1.74", "6.81", "", "TRI2", "15", "1.66", "1.71", "10.21"],
            ["EUSK", "15", "2.41", "2.06", "10.08", "", "TROS", "15", "3.23", "2.18", "10.97"],
            ["GRAS", "15", "1.81", "1.25", "3.95", "", "WALF", "15", "1.99", "2.29", "14.42"],
            ["GRAZ", "15", "1.61", "1.39", "5.85", "", "WARN", "15", "1.36", "2.15", "10.25"],
            ["HERS", "15", "2.52", "1.58", "4.59", "", "WSRT", "15", "1.95", "1.48", "7.54"],
            ["JOZE", "15", "1.19", "1.03", "4.08", "", "WTZR", "15", "0.92", "1.29", "7.82"],
            ["KBG2", "15", "1.90", "2.00", "9.99", "", "YEBE", "15", "1.67", "1.11", "3.88"],
            ["LONG", "13", "1.53", "1.84", "4.24", "", "ZIM2", "15", "1.48", "1.59", "4.87"],
        ],
        ["NETWORK MEAN (38 stations):  E: 2.02 mm  |  N: 1.74 mm  |  U: 7.38 mm",
         "Best: THNV (U=3.12mm), YEBE (U=3.88mm)  |  Worst: WALF (U=14.42mm), ONSA (U=11.46mm)"]
    )

    # ========== SLIDE 18: Summary and Best Practices ==========
    add_content_slide(prs, "Summary and Best Practices", [
        "Multi-GNSS Processing: GPS + Galileo for AR; GLONASS/BeiDou for geometry",
        "Critical Products: High-rate clocks (30s), OSB phase biases, IGS20 antennas",
        "Model Fidelity: VMF3 troposphere, FES2014b ocean loading, consistent IGS20 frame",
        "Quality Monitoring: Daily fixing rates >90%, sub-mm horizontal repeatability",
        "Automated QC: Monitor satellite residuals, OTL signatures, convergence times",
        "Future: Expanded Galileo constellation, GLONASS CDMA (K2 satellites)"
    ])

    # Save the presentation
    output_path = "/home/ahunegnaw/Python_IGNSS/i-GNSS/pygnss_rt/frontend/Advanced_PPP-AR_Strategies_Bernese54.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == '__main__':
    create_presentation()
